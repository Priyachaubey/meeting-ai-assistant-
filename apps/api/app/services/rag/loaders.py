import io


class DocumentLoadError(Exception):
    """Raised for unsupported file types or files that fail to parse — never silently
    returns empty/placeholder text for a file that couldn't actually be read."""


def load_text(filename: str, content: bytes) -> str:
    """Dispatches on file extension to a real extractor. Raises DocumentLoadError for
    anything unsupported rather than guessing or returning an empty string — an empty
    "successful" ingest would silently produce a knowledge base entry with nothing in it,
    which is worse than a clear failure at upload time."""
    suffix = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if suffix == "pdf":
        return _load_pdf(content)
    if suffix == "docx":
        return _load_docx(content)
    if suffix in ("xlsx", "xls"):
        return _load_excel(content)
    if suffix == "pptx":
        return _load_pptx(content)
    if suffix in ("txt", "md", "markdown", "csv"):
        return _load_plain_text(content)

    raise DocumentLoadError(
        f"Unsupported file type '.{suffix}' for '{filename}'. Supported: pdf, docx, xlsx, xls, pptx, txt, md, csv."
    )


def _load_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
    except Exception as exc:
        raise DocumentLoadError(f"Could not parse PDF: {exc}") from exc
    text = "\n\n".join(p for p in pages if p.strip())
    if not text.strip():
        raise DocumentLoadError("PDF parsed but contained no extractable text (likely a scanned/image-only PDF — OCR isn't wired in).")
    return text


def _load_excel(content: bytes) -> str:
    """openpyxl, not pandas — pandas would work too but pulls in numpy as a real dependency
    for what's fundamentally just "read cell values out of sheets," which openpyxl does
    directly. Renders each row as pipe-separated values (matching the DOCX table-loading
    convention above) rather than trying to preserve exact spreadsheet layout, which RAG
    chunking doesn't need and would just add noise to."""
    from openpyxl import load_workbook

    try:
        workbook = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
        sheets = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    except Exception as exc:
        raise DocumentLoadError(f"Could not parse Excel file: {exc}") from exc
    text = "\n\n".join(sheets)
    if not text.strip():
        raise DocumentLoadError("Excel file parsed but contained no extractable text (all sheets empty).")
    return text


def _load_docx(content: bytes) -> str:
    from docx import Document

    try:
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                paragraphs.append(" | ".join(cell.text for cell in row.cells))
    except Exception as exc:
        raise DocumentLoadError(f"Could not parse DOCX: {exc}") from exc
    text = "\n".join(p for p in paragraphs if p.strip())
    if not text.strip():
        raise DocumentLoadError("DOCX parsed but contained no extractable text.")
    return text


def _load_pptx(content: bytes) -> str:
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(content))
        chunks = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
            if texts:
                chunks.append(f"[Slide {i}]\n" + "\n".join(texts))
    except Exception as exc:
        raise DocumentLoadError(f"Could not parse PPTX: {exc}") from exc
    text = "\n\n".join(chunks)
    if not text.strip():
        raise DocumentLoadError("PPTX parsed but contained no extractable text.")
    return text


def _load_plain_text(content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        try:
            return content.decode("latin-1")
        except Exception as exc:
            raise DocumentLoadError(f"Could not decode text file: {exc}") from exc
